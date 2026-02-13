import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Read the data
base_dir = Path("Universe Data")
file_path = max(
    base_dir.glob("P123_Screen_*.csv"),
    key=lambda p: p.stat().st_mtime
)

# Read metadata
with file_path.open("r", encoding="utf-8-sig") as f:
    lines = [line.rstrip("\n") for line in f]

def clean_metadata_line(line):
    parts = [p.strip() for p in line.split(",")]
    for p in parts:
        if p:
            return p
    return None

metadata = {
    "title": clean_metadata_line(lines[0]) if len(lines) > 0 else None,
    "date": clean_metadata_line(lines[1]) if len(lines) > 1 else None,
    "notes": clean_metadata_line(lines[2]) if len(lines) > 2 else None,
}

# Read the dataframe
df = pd.read_csv(file_path, skiprows=3)
df = df.loc[:, ~df.columns.str.startswith("Unnamed")]

# Convert MktCap to numeric (remove commas and convert to float)
df['MktCap'] = pd.to_numeric(df['MktCap'].astype(str).str.replace(',', ''), errors='coerce')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(2)

title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = metadata["title"]
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add date
date_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(0.5))
date_frame = date_box.text_frame
date_frame.text = f"Analysis Date: {metadata['date']}"
date_frame.paragraphs[0].font.size = Pt(20)
date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 2: Summary Statistics
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Universe Summary Statistics"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Calculate statistics
total_stocks = len(df)
avg_mktcap = df['MktCap'].mean()
median_mktcap = df['MktCap'].median()
total_mktcap = df['MktCap'].sum()
num_sectors = df['SectorCode'].nunique()

# Add statistics text box
stats_text = f"""
Total Number of Stocks: {total_stocks:,}

Total Market Cap: ${total_mktcap:,.2f}M

Average Market Cap: ${avg_mktcap:,.2f}M

Median Market Cap: ${median_mktcap:,.2f}M

Number of Sectors: {num_sectors}
"""

stats_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
stats_frame = stats_box.text_frame
stats_frame.text = stats_text.strip()
for paragraph in stats_frame.paragraphs:
    paragraph.font.size = Pt(24)
    paragraph.space_after = Pt(12)

# Slide 3: Sector Breakdown
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Sector Breakdown"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Sector breakdown
sector_counts = df['SectorCode'].value_counts()
sector_text = "Stock Count by Sector:\n\n"
for sector, count in sector_counts.items():
    pct = (count / total_stocks) * 100
    sector_text += f"{sector}: {count:,} ({pct:.1f}%)\n"

sector_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
sector_frame = sector_box.text_frame
sector_frame.text = sector_text
for paragraph in sector_frame.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 4: Top 15 Stocks by Market Cap
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Top 15 Stocks by Market Cap"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Get top 15 stocks
top_stocks = df.nlargest(15, 'MktCap')[['Ticker', 'Name', 'MktCap', 'SectorCode']]

# Add table
rows = len(top_stocks) + 1
cols = 4
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(1.2)
table.columns[1].width = Inches(4)
table.columns[2].width = Inches(2)
table.columns[3].width = Inches(1.8)

# Header row
headers = ['Ticker', 'Name', 'Market Cap ($M)', 'Sector']
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
for row_idx, (_, row_data) in enumerate(top_stocks.iterrows(), start=1):
    table.cell(row_idx, 0).text = str(row_data['Ticker'])
    table.cell(row_idx, 1).text = str(row_data['Name'])[:40]  # Truncate long names
    table.cell(row_idx, 2).text = f"${row_data['MktCap']:,.2f}"
    table.cell(row_idx, 3).text = str(row_data['SectorCode'])

    # Set font size for data rows
    for col_idx in range(cols):
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)

# Save presentation
output_path = Path("results.pptx")
prs.save(output_path)
print(f"PowerPoint presentation created successfully: {output_path.absolute()}")
